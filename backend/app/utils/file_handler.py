import os
import hashlib
import aiofiles
import asyncio
import sys
from typing import List

from langchain_core.documents import Document

from app.core.logger_handler import logger
from app.utils.path_tool import get_abstract_path


class FontBBoxStreamFilter:
    def __init__(self, stream):
        self.stream = stream
        
    def write(self, data):
        if isinstance(data, bytes):
            if b'FontBBox from font descriptor' not in data:
                self.stream.write(data)
        else:
            if 'FontBBox from font descriptor' not in data:
                self.stream.write(data)
            
    def flush(self):
        self.stream.flush()

sys.stderr = FontBBoxStreamFilter(sys.stderr)


# ── MD5 计算 ──

async def get_file_md5_hex(file_path: str) -> str:
    """获取文件的md5值"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    if not os.path.exists(abs_file_path):
        logger.error(f"【md5计算】文件路径 {abs_file_path} 不存在")
        return ""
    if not os.path.isfile(abs_file_path):
        logger.error(f"【md5计算】文件路径 {abs_file_path} 不是文件")
        return ""
    md5_object = hashlib.md5()
    chunk_size = 1024
    try:
        async with aiofiles.open(abs_file_path, "rb") as f:
            while chunk := await f.read(chunk_size):
                md5_object.update(chunk)
    except Exception as e:
        logger.error(f"【md5计算】读取文件 {abs_file_path} 时出错: {e}")
        return ""
    return md5_object.hexdigest()

def get_file_md5_hex_sync(file_path: str) -> str:
    """同步获取文件的md5值（用于多线程场景）"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    if not os.path.exists(abs_file_path):
        logger.error(f"【md5计算】文件路径 {abs_file_path} 不存在")
        return ""
    if not os.path.isfile(abs_file_path):
        logger.error(f"【md5计算】文件路径 {abs_file_path} 不是文件")
        return ""
    md5_object = hashlib.md5()
    chunk_size = 1024
    try:
        with open(abs_file_path, "rb") as f:
            while chunk := f.read(chunk_size):
                md5_object.update(chunk)
    except Exception as e:
        logger.error(f"【md5计算】读取文件 {abs_file_path} 时出错: {e}")
        return ""
    return md5_object.hexdigest()


# ── 目录遍历 ──

async def listdir_allowed_type(path: str, allowed_types: tuple[str]) -> tuple:
    """获取指定目录下所有允许的文件类型"""
    abs_path = get_abstract_path(path) if not os.path.isabs(path) else path
    if not os.path.exists(abs_path):
        logger.error(f"【文件列表】目录路径 {abs_path} 不存在")
        return ()
    if not os.path.isdir(abs_path):
        logger.error(f"【文件列表】目录路径 {abs_path} 不是目录")
        return ()
    file_list = []
    for f in await asyncio.to_thread(os.listdir, abs_path):
        if f.endswith(allowed_types):
            file_path = os.path.join(abs_path, f)
            file_list.append(file_path)
    return tuple(file_list)


# ── 文档加载器 ──

async def pdf_loader(file_path: str, password: str = None) -> List[Document]:
    """PDF 加载器：使用 pypdf 提取文本（无系统依赖）"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path

    import pypdf
    reader = pypdf.PdfReader(abs_file_path)
    if password and reader.is_encrypted:
        reader.decrypt(password)
    docs = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        docs.append(Document(page_content=text, metadata={"page": i + 1, "source": abs_file_path}))
    return docs


async def txt_loader(file_path: str) -> List[Document]:
    """TXT 加载器"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    encodings = ['utf-8', 'gbk']
    for encoding in encodings:
        try:
            async with aiofiles.open(abs_file_path, 'r', encoding=encoding) as f:
                content = await f.read()
            return [Document(page_content=content, metadata={"source": abs_file_path})]
        except Exception as e:
            logger.warning(f"【文本文件加载】使用编码 {encoding} 失败: {e}")
            continue
    return []


async def word_loader(file_path: str) -> List[Document]:
    """DOCX 加载器"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    try:
        from unstructured.partition.docx import partition_docx
        from app.rag.document_handler.format_preserver import aggregate_by_length
        elements = partition_docx(filename=abs_file_path)
        return aggregate_by_length(elements, abs_file_path)
    except Exception as e:
        logger.error(f"【WORD文件加载】失败: {e}")
        return []


async def markdown_loader(file_path: str) -> List[Document]:
    """Markdown 加载器"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    try:
        from unstructured.partition.md import partition_md
        from app.rag.document_handler.format_preserver import aggregate_by_length
        elements = partition_md(filename=abs_file_path)
        return aggregate_by_length(elements, abs_file_path)
    except Exception as e:
        logger.error(f"【Markdown文件加载】失败: {e}")
        return []


async def ppt_loader(file_path: str) -> List[Document]:
    """PPTX 加载器"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    try:
        from unstructured.partition.pptx import partition_pptx
        from app.rag.document_handler.format_preserver import aggregate_by_slide
        elements = partition_pptx(filename=abs_file_path)
        return aggregate_by_slide(elements, abs_file_path)
    except Exception as e:
        logger.error(f"【PPT文件加载】unstructured 失败，尝试 fallback: {e}", exc_info=True)
        return _pptx_fallback_load(abs_file_path)


# ── 同步版本（用于多线程场景）──

def pdf_loader_sync(file_path: str, password: str = None) -> List[Document]:
    """同步加载PDF文件内容（使用 pypdf，无系统依赖）"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path

    import pypdf
    reader = pypdf.PdfReader(abs_file_path)
    if password and reader.is_encrypted:
        reader.decrypt(password)
    docs = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        docs.append(Document(page_content=text, metadata={"page": i + 1, "source": abs_file_path}))
    return docs


def txt_loader_sync(file_path: str) -> List[Document]:
    """同步加载TXT文件内容"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    encodings = ['utf-8', 'gbk']
    for encoding in encodings:
        try:
            with open(abs_file_path, 'r', encoding=encoding) as f:
                content = f.read()
            return [Document(page_content=content, metadata={"source": abs_file_path})]
        except Exception:
            continue
    return []


def word_loader_sync(file_path: str) -> List[Document]:
    """同步加载DOCX文件内容"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    try:
        from unstructured.partition.docx import partition_docx
        from app.rag.document_handler.format_preserver import aggregate_by_length
        elements = partition_docx(filename=abs_file_path)
        if elements:
            return aggregate_by_length(elements, abs_file_path)

        # partition_docx 返回空 → 尝试回退提取浮动文本框中的内容
        # （Aspose.Words 等自动化工具生成的 docx 使用 wp:anchor 浮动形状）
        logger.debug(f"【WORD文件加载(同步)】partition_docx 返回空，尝试回退: {file_path}")
        return _extract_textbox_docx(abs_file_path)
    except Exception as e:
        logger.error(f"【WORD文件加载(同步)】失败: {e}")
        return []


def _extract_textbox_docx(file_path: str) -> List[Document]:
    """从 docx 浮动文本框 (wps:wsp / wp:anchor) 中手动提取文本

    方案 A 精神：扩展 XPath 覆盖 wp:anchor 浮动形状
    方案 B 精神：直接遍历 w:txbxContent 提取全部文本内容
    提取结果复用 aggregate_by_length 进行聚合。
    """
    from lxml import etree
    from docx import Document as DocxDocument
    from app.rag.document_handler.format_preserver import aggregate_by_length

    W_NS = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'

    doc = DocxDocument(file_path)
    root = doc.element.getroottree().getroot()

    # 遍历所有 w:txbxContent（文本框内容容器），收集每个内部段落的文字
    items: list[str] = []
    for txbx in root.iter(f'{{{W_NS}}}txbxContent'):
        for p in txbx.iter(f'{{{W_NS}}}p'):
            text = ''.join(t.text or '' for t in p.iter(f'{{{W_NS}}}t'))
            if text.strip():
                items.append(text.strip())

    if not items:
        return []

    # 构建简单对象供 aggregate_by_length 使用（category=None → 纯文本，无 Markdown 标记）
    class _TextElement:
        __slots__ = ('_text',)
        def __init__(self, text: str):
            self._text = text
        def __str__(self) -> str:
            return self._text
        @property
        def category(self):
            return None

    return aggregate_by_length([_TextElement(t) for t in items], file_path)


def markdown_loader_sync(file_path: str) -> List[Document]:
    """同步加载Markdown文件内容"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    try:
        from unstructured.partition.md import partition_md
        from app.rag.document_handler.format_preserver import aggregate_by_length
        elements = partition_md(filename=abs_file_path)
        return aggregate_by_length(elements, abs_file_path)
    except Exception as e:
        logger.error(f"【Markdown文件加载(同步)】失败: {e}")
        return []
    except Exception as e:
        logger.error(f"【Markdown文件加载(同步)】失败: {e}")
        return []


def ppt_loader_sync(file_path: str) -> List[Document]:
    """同步加载PPTX文件内容"""
    abs_file_path = get_abstract_path(file_path) if not os.path.isabs(file_path) else file_path
    try:
        from unstructured.partition.pptx import partition_pptx
        from app.rag.document_handler.format_preserver import aggregate_by_slide
        elements = partition_pptx(filename=abs_file_path)
        return aggregate_by_slide(elements, abs_file_path)
    except Exception as e:
        logger.error(f"【PPT文件加载(同步)】unstructured 失败，尝试 fallback: {e}", exc_info=True)
        return _pptx_fallback_load(abs_file_path)


def _pptx_fallback_load(file_path: str) -> List[Document]:
    """PPTX fallback 加载器：直接用 python-pptx 解析，绕过 unstructured 的 spacy 依赖

    当 unstructured 的 partition_pptx 因 spacy 模型未安装/网络超时失败时启用。
    不做文本类型分类（Title/NarrativeText 判断），但仍能提取所有文本，
    并通过 aggregate_by_slide 完成幻灯片聚合与元信息注入。
    """
    try:
        from pptx import Presentation
        from app.rag.document_handler.format_preserver import aggregate_by_slide

        class _FakeElement:
            """模拟 unstructured element，供 aggregate_by_slide 使用"""
            __slots__ = ('_text', 'category', 'metadata')

            def __init__(self, text: str, category: str, page_number: int):
                self._text = text
                self.category = category
                self.metadata = type('_M', (), {'page_number': page_number})()

            def __str__(self) -> str:
                return self._text

        prs = Presentation(file_path)
        elements = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            title_shape = slide.shapes.title
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    # 简单分类：标题 shape → Title，有缩进 → ListItem，其余 → NarrativeText
                    if title_shape is not None and shape == title_shape:
                        category = "Title"
                    elif para.level > 0:
                        category = "ListItem"
                    else:
                        category = "NarrativeText"
                    elements.append(_FakeElement(text, category, slide_idx))

        logger.info(f"【PPT fallback】python-pptx 提取 {len(elements)} 个元素")
        return aggregate_by_slide(elements, file_path)
    except Exception as e:
        logger.error(f"【PPT fallback】也失败: {e}", exc_info=True)
        return []
